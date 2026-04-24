# PacketArch — OT Traffic Simulation Platform
# Copyright (c) 2026 Rocky Smith <rocky.d.smith@proton.me>
# Licensed under GPL-3.0. See LICENSE at the repo root.
"""Generate the executive architecture review PowerPoint for PacketArch."""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent.parent / "docs" / "PacketArch_Architecture_Review.pptx"

# --- Brand palette ----------------------------------------------------------
NAVY = RGBColor(0x0B, 0x1F, 0x3A)
TEAL = RGBColor(0x00, 0xA6, 0xA6)
SLATE = RGBColor(0x33, 0x3F, 0x4E)
LIGHT = RGBColor(0xF4, 0xF6, 0xF8)
ACCENT = RGBColor(0xE8, 0x8B, 0x00)
MUTED = RGBColor(0x6B, 0x73, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _add_rect(slide, left, top, width, height, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def _add_text(
    slide, left, top, width, height, text, *,
    size=14, bold=False, color=SLATE, align=PP_ALIGN.LEFT, font="Calibri",
):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = font
    run.font.color.rgb = color
    return tb


def _bullets(slide, left, top, width, height, items, *, size=14, color=SLATE, bold_lead=True):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(4)
        if isinstance(item, tuple):
            lead, rest = item
            r1 = p.add_run()
            r1.text = f"• {lead}"
            r1.font.size = Pt(size)
            r1.font.name = "Calibri"
            r1.font.bold = bold_lead
            r1.font.color.rgb = NAVY
            r2 = p.add_run()
            r2.text = f"  {rest}"
            r2.font.size = Pt(size)
            r2.font.name = "Calibri"
            r2.font.color.rgb = color
        else:
            r = p.add_run()
            r.text = f"• {item}"
            r.font.size = Pt(size)
            r.font.name = "Calibri"
            r.font.color.rgb = color
    return tb


def _slide(prs, layout_idx=6):  # blank layout
    s = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    # background
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()
    return s


def _header(slide, title, subtitle=None, number=None):
    # left accent bar
    _add_rect(slide, 0, 0, Inches(0.3), SLIDE_H, NAVY)
    # title bar
    _add_rect(slide, Inches(0.3), 0, SLIDE_W - Inches(0.3), Inches(0.95), LIGHT)
    _add_text(slide, Inches(0.55), Inches(0.18), Inches(11), Inches(0.55),
              title, size=26, bold=True, color=NAVY)
    if subtitle:
        _add_text(slide, Inches(0.55), Inches(0.60), Inches(11), Inches(0.32),
                  subtitle, size=12, color=MUTED)
    if number:
        _add_text(slide, Inches(12.2), Inches(0.25), Inches(1.0), Inches(0.4),
                  number, size=11, bold=True, color=TEAL, align=PP_ALIGN.RIGHT)


def _footer(slide, label):
    _add_text(slide, Inches(0.55), Inches(7.15), Inches(8), Inches(0.3),
              "PacketArch — Architecture Review", size=9, color=MUTED)
    _add_text(slide, Inches(9), Inches(7.15), Inches(4), Inches(0.3),
              label, size=9, color=MUTED, align=PP_ALIGN.RIGHT)


# ============================================================================
# SLIDES
# ============================================================================

def slide_title(prs):
    s = _slide(prs)
    _add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    _add_rect(s, 0, Inches(3.2), SLIDE_W, Inches(0.05), TEAL)
    _add_text(s, Inches(0.75), Inches(2.3), Inches(12), Inches(0.8),
              "PacketArch", size=52, bold=True, color=WHITE)
    _add_text(s, Inches(0.75), Inches(3.35), Inches(12), Inches(0.6),
              "Architecture & Deployment Review", size=28, color=TEAL)
    _add_text(s, Inches(0.75), Inches(4.1), Inches(12), Inches(0.4),
              "OT Traffic Simulation Platform — Executive Walk-through", size=16, color=LIGHT)
    _add_text(s, Inches(0.75), Inches(6.8), Inches(12), Inches(0.3),
              "github.com/ip-aegis/PacketArch", size=11, color=MUTED)


def slide_overview(prs):
    s = _slide(prs)
    _header(s, "What PacketArch Is", "A one-page orientation", "01")
    _add_text(s, Inches(0.55), Inches(1.2), Inches(12), Inches(0.5),
              "Generate realistic OT / ICS network traffic on demand — for training, testing, and validation of security tools such as Cisco Cyber Vision.",
              size=15, color=SLATE)
    # three panels
    cards = [
        ("Design", "Drag-and-drop scenario studio: zones, devices, protocol flows on an xyflow canvas.", TEAL),
        ("Generate", "22+ OT protocols, 295 vendor fingerprints, process simulation, adaptive timing, attack playbooks.", ACCENT),
        ("Deliver", "PCAP export for offline use, or live injection via remote agents phoning home over WebSocket.", NAVY),
    ]
    x = Inches(0.55)
    for title, body, color in cards:
        _add_rect(s, x, Inches(2.2), Inches(4.05), Inches(3.4), LIGHT)
        _add_rect(s, x, Inches(2.2), Inches(4.05), Inches(0.55), color)
        _add_text(s, x + Inches(0.2), Inches(2.28), Inches(3.7), Inches(0.4),
                  title, size=18, bold=True, color=WHITE)
        _add_text(s, x + Inches(0.2), Inches(2.9), Inches(3.7), Inches(2.6),
                  body, size=13, color=SLATE)
        x += Inches(4.25)
    _add_text(s, Inches(0.55), Inches(6.1), Inches(12), Inches(0.5),
              "Used for: security tool validation, OT SOC analyst training, procurement demos, compliance testing (IEC 62443 conduits).",
              size=12, color=MUTED)
    _footer(s, "Overview")


def slide_system_arch(prs):
    s = _slide(prs)
    _header(s, "System Architecture", "Three tiers + remote data plane", "02")

    # Browser
    _add_rect(s, Inches(0.6), Inches(1.35), Inches(12.1), Inches(0.8), LIGHT)
    _add_text(s, Inches(0.8), Inches(1.55), Inches(11.7), Inches(0.4),
              "Browser — React 19 SPA (Vite, TypeScript, Ant Design, xyflow canvas, Zustand, TanStack Query)",
              size=13, bold=True, color=NAVY)

    # Edge / nginx
    _add_rect(s, Inches(0.6), Inches(2.3), Inches(12.1), Inches(0.7), SLATE)
    _add_text(s, Inches(0.8), Inches(2.48), Inches(11.7), Inches(0.4),
              "Edge — Nginx (HTTPS / HTTP/2, self-signed TLS, WebSocket proxy, SPA fallback)",
              size=13, bold=True, color=WHITE)

    # App tier
    _add_rect(s, Inches(0.6), Inches(3.15), Inches(5.9), Inches(2.6), NAVY)
    _add_text(s, Inches(0.8), Inches(3.28), Inches(5.5), Inches(0.4),
              "Application", size=14, bold=True, color=TEAL)
    _bullets(s, Inches(0.8), Inches(3.68), Inches(5.5), Inches(2.0), [
        "FastAPI (async) — 25+ routers",
        "Celery worker — PCAP & deployment jobs",
        "MCP server — JSON-RPC + SSE for Claude",
        "Agent hub — WebSocket /ws/agent",
        "Protocol engines (shared with agent)",
    ], size=12, color=WHITE, bold_lead=False)

    # Data tier
    _add_rect(s, Inches(6.7), Inches(3.15), Inches(6.0), Inches(2.6), TEAL)
    _add_text(s, Inches(6.9), Inches(3.28), Inches(5.5), Inches(0.4),
              "Data", size=14, bold=True, color=WHITE)
    _bullets(s, Inches(6.9), Inches(3.68), Inches(5.5), Inches(2.0), [
        "PostgreSQL + TimescaleDB 15 (async, SQLAlchemy 2.0, Alembic)",
        "Redis 7 — Celery broker + cache",
        "Volumes: pcap_uploads, pcap_output, ssl_certs",
        "Secrets: .env (server-init.sh), JWT HS256",
    ], size=12, color=WHITE, bold_lead=False)

    # Remote agents
    _add_rect(s, Inches(0.6), Inches(5.9), Inches(12.1), Inches(1.1), ACCENT)
    _add_text(s, Inches(0.8), Inches(6.02), Inches(11.7), Inches(0.4),
              "Remote Traffic Agents", size=14, bold=True, color=WHITE)
    _add_text(s, Inches(0.8), Inches(6.42), Inches(11.7), Inches(0.5),
              "Python 3.12 container, network_mode=host, CAP_NET_ADMIN/RAW. Phones home over WebSocket — no inbound ports. Ships with shared protocol_engines.",
              size=12, color=WHITE)

    _footer(s, "System Architecture")


def slide_tech_stack(prs):
    s = _slide(prs)
    _header(s, "Technology Stack", "Versions pinned in pyproject.toml / package.json", "03")

    groups = [
        ("Backend", NAVY, [
            "Python 3.11 · FastAPI 0.109 · Uvicorn 0.27",
            "SQLAlchemy 2.0 async · asyncpg · Alembic",
            "Celery 5.3 · Redis 7",
            "Pydantic 2.5 · python-statemachine 2.1",
            "Scapy 2.5 · NumPy / SciPy (process sim)",
            "Anthropic SDK 0.79 · httpx · ldap3",
        ]),
        ("Frontend", TEAL, [
            "React 19.2 · TypeScript 5.9 (strict)",
            "Vite 7.2 · React Router 7.6",
            "Ant Design 5.24 · @xyflow/react 12.6",
            "Zustand 5 · TanStack Query 5.72 · Axios",
            "Recharts · react-markdown · Zod",
            "ESLint 9 · Vitest 4 · MSW 2",
        ]),
        ("Platform & Ops", ACCENT, [
            "Docker + docker-compose 2.32",
            "Nginx alpine (TLS termination, HTTP/2)",
            "PostgreSQL 15 + TimescaleDB",
            "Poetry · pnpm · pre-commit",
            "Ruff · Black · mypy · pytest-asyncio",
            "GitHub Actions CI + rsync deploy",
        ]),
    ]

    x = Inches(0.55)
    for title, color, items in groups:
        _add_rect(s, x, Inches(1.3), Inches(4.1), Inches(5.5), LIGHT)
        _add_rect(s, x, Inches(1.3), Inches(4.1), Inches(0.55), color)
        _add_text(s, x + Inches(0.2), Inches(1.38), Inches(3.8), Inches(0.4),
                  title, size=16, bold=True, color=WHITE)
        tb = _bullets(s, x + Inches(0.2), Inches(2.0), Inches(3.8), Inches(4.7),
                      items, size=12, color=SLATE, bold_lead=False)
        x += Inches(4.25)

    _add_text(s, Inches(0.55), Inches(6.95), Inches(12), Inches(0.3),
              "Single-language per tier, modern versions across the board; no legacy Python 2, no AngularJS, no jQuery.",
              size=11, color=MUTED)
    _footer(s, "Technology Stack")


def slide_backend(prs):
    s = _slide(prs)
    _header(s, "Backend — FastAPI Monolith", "app/main.py registers 25+ routers + lifespan hooks", "04")

    # Left column: subsystems
    _add_text(s, Inches(0.55), Inches(1.25), Inches(6), Inches(0.4),
              "Subsystems", size=16, bold=True, color=NAVY)
    _bullets(s, Inches(0.55), Inches(1.7), Inches(6), Inches(5.3), [
        ("api/routes/", "REST endpoints: scenarios, deployments, agents, attacks, adaptation, AI, CV, LDAP"),
        ("services/", "~32 services: agent_manager, scenario_remediation, cyber_vision, LDAP, health_monitor"),
        ("protocol_engines/", "22+ engines + identity builders + adaptive/attacks/process_sim/ambient"),
        ("mcp_server/", "JSON-RPC 2.0 + HTTP SSE; ~6.3K LOC of tool implementations for Claude"),
        ("ai_services/", "Chat, scenario generation, preview validation; MCP tools exposed to Claude"),
        ("models/ + alembic/", "14 async SQLAlchemy models, migrations via Alembic"),
        ("traffic_generator/", "EventScheduler, PCAP writer, Celery tasks"),
    ], size=12)

    # Right column: highlights
    _add_text(s, Inches(6.9), Inches(1.25), Inches(6), Inches(0.4),
              "Notable Patterns", size=16, bold=True, color=NAVY)
    _bullets(s, Inches(6.9), Inches(1.7), Inches(6), Inches(5.3), [
        ("Async all the way:", "FastAPI + asyncpg + httpx; Celery handles CPU/IO-heavy PCAP jobs"),
        ("Shared engines:", "Agent Docker build stages protocol_engines from backend — single source of truth"),
        ("Unified orchestrator:", "Same code drives PCAP generation and live agent injection"),
        ("Typed end-to-end:", "Pydantic 2 on requests/responses; mypy strict; custom PacketArchError hierarchy"),
        ("Cache layer:", "FingerprintCache singleton; batched seed via db.add_all(); .in_() for bulk lookups"),
    ], size=12)

    _footer(s, "Backend")


def slide_frontend(prs):
    s = _slide(prs)
    _header(s, "Frontend — Scenario Studio SPA", "React 19 + Vite + @xyflow/react canvas", "05")

    _add_text(s, Inches(0.55), Inches(1.25), Inches(6), Inches(0.4),
              "14 Primary Pages", size=16, bold=True, color=NAVY)
    _bullets(s, Inches(0.55), Inches(1.7), Inches(6), Inches(5.3), [
        "ScenarioStudioPage — xyflow canvas (DeviceNode, ZoneNode, FlowEdge)",
        "AIScenarioWizardPage — natural language → scenario via Claude + MCP",
        "GuidedBuilderPage — step-by-step scenario construction",
        "DeploymentsPage — live deployment lifecycle, agent assignment",
        "LiveTrafficDashboardPage — real-time bytes/packets per protocol",
        "CyberVisionPage — device comparison + MAC/IP matching",
        "FingerprintingLibraryPage · CVEBrowserPage · IPManagementPage",
        "SettingsPage — CV, LDAP, agent image builder (admin)",
    ], size=12)

    _add_text(s, Inches(6.9), Inches(1.25), Inches(6), Inches(0.4),
              "State & Data", size=16, bold=True, color=NAVY)
    _bullets(s, Inches(6.9), Inches(1.7), Inches(6), Inches(5.3), [
        ("15 Zustand stores:", "auth, scenario, ui, agents, liveDashboard, deployment, attack, history…"),
        ("TanStack Query:", "Server state cache, optimistic mutations"),
        ("Axios client:", "JWT auto-injection + typed error extraction"),
        ("Protocol types:", "Discriminated unions with type guards for 22 protocols"),
        ("Ant Design dark theme:", "Consistent component library, form validation via react-hook-form + Zod"),
        ("Testing:", "Vitest + jsdom + MSW for API mocking"),
    ], size=12)

    _footer(s, "Frontend")


def slide_agent(prs):
    s = _slide(prs)
    _header(s, "Remote Traffic Agent", "Phone-home WebSocket model, no inbound ports", "06")

    # Diagram
    _add_rect(s, Inches(0.55), Inches(1.3), Inches(12.2), Inches(2.1), LIGHT)
    _add_rect(s, Inches(1.0), Inches(1.8), Inches(2.8), Inches(1.1), NAVY)
    _add_text(s, Inches(1.0), Inches(2.1), Inches(2.8), Inches(0.5),
              "Customer Site\nPacketArch Agent", size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # arrow
    arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                               Inches(4.1), Inches(2.15), Inches(5.1), Inches(0.45))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = TEAL
    arrow.line.fill.background()
    _add_text(s, Inches(4.2), Inches(1.75), Inches(4.9), Inches(0.4),
              "wss:// /ws/agent?token=…  (outbound 443)", size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    _add_rect(s, Inches(9.5), Inches(1.8), Inches(2.8), Inches(1.1), TEAL)
    _add_text(s, Inches(9.5), Inches(2.1), Inches(2.8), Inches(0.5),
              "PacketArch Server\nAgent Hub", size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Details columns
    _add_text(s, Inches(0.55), Inches(3.6), Inches(6), Inches(0.4),
              "Container Profile", size=14, bold=True, color=NAVY)
    _bullets(s, Inches(0.55), Inches(4.0), Inches(6), Inches(3.0), [
        "Python 3.12-slim + libpcap + Docker CLI",
        "network_mode: host, CAP_NET_ADMIN / NET_RAW for raw injection",
        "Docker socket mounted → agent self-updates via UPDATE_AGENT",
        "Current version: 1.26.0 (semver gated on protocol_engines changes)",
        "Installed by curl-piped /agent/install.sh from server",
    ], size=12)

    _add_text(s, Inches(6.9), Inches(3.6), Inches(6), Inches(0.4),
              "Protocol Summary", size=14, bold=True, color=NAVY)
    _bullets(s, Inches(6.9), Inches(4.0), Inches(6), Inches(3.0), [
        ("Server → Agent:", "START / STOP / UPDATE_SCENARIO, ADAPT_TRAFFIC, START/STOP_ATTACK, UPDATE_AGENT, PING"),
        ("Agent → Server:", "STATUS, INTERFACES, HEARTBEAT (CPU/mem/version), ERROR, UPDATE_STATUS"),
        ("Shared code:", "Docker build stages backend/app/protocol_engines into agent image"),
    ], size=12)

    _footer(s, "Remote Agent")


def slide_traffic_engine(prs):
    s = _slide(prs)
    _header(s, "Traffic Generation Engine", "UnifiedOrchestrator composes five peers", "07")

    # center orchestrator
    _add_rect(s, Inches(5.1), Inches(3.4), Inches(3.1), Inches(1.3), NAVY)
    _add_text(s, Inches(5.1), Inches(3.6), Inches(3.1), Inches(0.45),
              "UnifiedOrchestrator", size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _add_text(s, Inches(5.1), Inches(4.05), Inches(3.1), Inches(0.55),
              "timed (PCAP) or perpetual (live)", size=11, color=TEAL, align=PP_ALIGN.CENTER)

    peers = [
        ("Protocol Engines", "State machines per flow\n(Modbus, S7, EIP, BACnet…)", Inches(0.55), Inches(1.3), TEAL),
        ("Adaptive Controller", "Micro-variations, drift,\ntime-of-day shaping", Inches(9.7), Inches(1.3), TEAL),
        ("Attack Orchestrator", "6 playbooks, 15+ Snort/\nSuricata-triggering actions", Inches(0.55), Inches(5.2), ACCENT),
        ("Process Simulation", "ODE models for sensors:\nwater, power, oil/gas, mfg", Inches(9.7), Inches(5.2), ACCENT),
    ]
    for title, body, left, top, color in peers:
        _add_rect(s, left, top, Inches(3.1), Inches(1.3), LIGHT)
        _add_rect(s, left, top, Inches(3.1), Inches(0.4), color)
        _add_text(s, left + Inches(0.1), top + Inches(0.03), Inches(2.9), Inches(0.35),
                  title, size=12, bold=True, color=WHITE)
        _add_text(s, left + Inches(0.1), top + Inches(0.45), Inches(2.9), Inches(0.85),
                  body, size=11, color=SLATE)

    # bottom band: ambient noise
    _add_rect(s, Inches(4.0), Inches(5.2), Inches(5.3), Inches(1.3), SLATE)
    _add_rect(s, Inches(4.0), Inches(5.2), Inches(5.3), Inches(0.4), NAVY)
    _add_text(s, Inches(4.1), Inches(5.23), Inches(5.1), Inches(0.35),
              "Ambient / Broadcast Ecosystem", size=12, bold=True, color=WHITE)
    _add_text(s, Inches(4.1), Inches(5.65), Inches(5.1), Inches(0.85),
              "LLDP · STP · CDP · DHCP · NTP · IGMP · BACnet Who-Is · PROFINET DCP · SNMP traps",
              size=11, color=WHITE)

    _add_text(s, Inches(0.55), Inches(6.75), Inches(12.2), Inches(0.35),
              "Output swap: PcapOutput(file) for offline PCAPs · LiveOutput(iface) for agent-side injection. Same generation code.",
              size=11, color=MUTED, align=PP_ALIGN.CENTER)
    _footer(s, "Traffic Engine")


def slide_protocols(prs):
    s = _slide(prs)
    _header(s, "Protocol Coverage", "22+ OT / ICS protocols, production-grade state machines", "08")

    protos = [
        ("Modbus TCP", "502", "Mfg, Utilities"),
        ("EtherNet/IP", "44818 / 2222", "Rockwell, ODVA"),
        ("PROFINET", "Layer 2", "Siemens, Mfg"),
        ("S7comm / S7+", "102", "Siemens PLCs"),
        ("BACnet/IP", "47808", "Building Auto"),
        ("SNMP / NTCIP", "161 / 162", "All vendors"),
        ("OPC UA", "4840", "Modern OT"),
        ("DNP3", "20000", "Utilities, SCADA"),
        ("IEC 60870-104", "2404", "European utilities"),
        ("IEC 61850", "102", "Substations, GOOSE"),
        ("EtherCAT", "Layer 2", "Motion control"),
        ("FINS", "9600", "Omron"),
        ("CODESYS", "11740", "Soft-PLCs"),
        ("SLMP", "2000", "Mitsubishi"),
        ("PCCC", "44818", "Legacy AB"),
        ("FANUC RoboGuide", "Custom", "Robotics"),
        ("CDP / LLDP", "Layer 2", "Discovery"),
        ("Cloud services", "443", "EWON, TeamViewer"),
    ]
    # grid 3 columns x 6 rows
    col_w = Inches(4.1)
    row_h = Inches(0.62)
    start_left = Inches(0.55)
    start_top = Inches(1.3)
    for i, (name, port, notes) in enumerate(protos):
        r, c = divmod(i, 3)
        left = start_left + c * col_w + Inches(0.15) * c
        top = start_top + r * row_h
        _add_rect(s, left, top, col_w, Inches(0.55), LIGHT)
        _add_rect(s, left, top, Inches(0.15), Inches(0.55), TEAL)
        _add_text(s, left + Inches(0.25), top + Inches(0.04), col_w - Inches(0.3), Inches(0.28),
                  name, size=12, bold=True, color=NAVY)
        _add_text(s, left + Inches(0.25), top + Inches(0.28), col_w - Inches(0.3), Inches(0.25),
                  f"{port}   ·   {notes}", size=10, color=MUTED)

    _add_text(s, Inches(0.55), Inches(6.7), Inches(12.2), Inches(0.5),
              "All engines extend a common base (startup → poll → shutdown) and plug into the identity builder system so each device responds as its real vendor/model.",
              size=11, color=MUTED, align=PP_ALIGN.CENTER)
    _footer(s, "Protocols")


def slide_fingerprints(prs):
    s = _slide(prs)
    _header(s, "Device Templates & Fingerprints", "Single source of truth in services/device_templates/", "09")

    # stats row
    stats = [
        ("295", "Device templates"),
        ("18", "Vendor modules"),
        ("IEEE", "Verified OUIs"),
        ("5", "Realism dimensions"),
    ]
    x = Inches(0.55)
    for n, label in stats:
        _add_rect(s, x, Inches(1.3), Inches(3.0), Inches(1.3), NAVY)
        _add_text(s, x, Inches(1.38), Inches(3.0), Inches(0.7),
                  n, size=34, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
        _add_text(s, x, Inches(2.05), Inches(3.0), Inches(0.5),
                  label, size=12, color=WHITE, align=PP_ALIGN.CENTER)
        x += Inches(3.1)

    _add_text(s, Inches(0.55), Inches(2.85), Inches(6), Inches(0.4),
              "What a Template Captures", size=14, bold=True, color=NAVY)
    _bullets(s, Inches(0.55), Inches(3.25), Inches(6), Inches(3.8), [
        "Network signatures (ports, protocols, banner strings)",
        "Protocol identities: sysName, station_name, CIP vendor_id, S7 SZL, BACnet ids",
        "Response timings (RTT, jitter bands) — per protocol",
        "Behavioral patterns (steady / warming-up / maintenance)",
        "Firmware versions derived deterministically from fingerprint",
        "Vendor-aligned MAC OUIs from protocol_engines/vendor_oui.py",
    ], size=12)

    _add_text(s, Inches(6.9), Inches(2.85), Inches(6), Inches(0.4),
              "Realism Guardrails (auto-enforced)", size=14, bold=True, color=NAVY)
    _bullets(s, Inches(6.9), Inches(3.25), Inches(6), Inches(3.8), [
        ("Naming:", "industrial-appropriate, unique (no device_001)"),
        ("Protocol accuracy:", "only vendor-supported protocols per device"),
        ("Completeness:", "every device in ≥1 flow; SNMP fallback if no partner"),
        ("Conduit compliance:", "cross-zone flows justified by IEC 62443"),
        ("MAC-vendor alignment:", "OUI must match declared vendor"),
    ], size=12)

    _footer(s, "Templates")


def slide_ai(prs):
    s = _slide(prs)
    _header(s, "AI & MCP Integration", "Claude drives scenario generation and assistance", "10")

    _add_text(s, Inches(0.55), Inches(1.25), Inches(12.2), Inches(0.4),
              "Model Context Protocol server exposes ~10 tool modules (6.3K LOC) to Anthropic Claude over HTTP + SSE.",
              size=13, color=SLATE)

    # two columns
    _add_text(s, Inches(0.55), Inches(1.9), Inches(6), Inches(0.4),
              "Tool Surface (Claude-callable)", size=14, bold=True, color=NAVY)
    _bullets(s, Inches(0.55), Inches(2.3), Inches(6), Inches(4.5), [
        "scenario_tools — CRUD scenarios",
        "device_tools / flow_tools — canvas edits",
        "protocol_tools — 1.5K LOC of protocol metadata",
        "fingerprint_tools — vendor/model lookup",
        "validation_tools — readiness + compliance",
        "layout_tools · addressing_tools · external_comm_tools",
        "ai_generation_tools — scenario preview + generation context",
    ], size=12)

    _add_text(s, Inches(6.9), Inches(1.9), Inches(6), Inches(0.4),
              "Where AI Shows Up", size=14, bold=True, color=NAVY)
    _bullets(s, Inches(6.9), Inches(2.3), Inches(6), Inches(4.5), [
        ("AI Scenario Wizard:", "NL prompt → fully-formed scenario JSON"),
        ("Scenario Preview:", "schema validation + remediation suggestions"),
        ("AI Assistant:", "context-aware chat sidebar, MCP tools in-loop"),
        ("Help system:", "embedded answers about PacketArch usage"),
        ("Sanitization layer:", "redacts secrets, caps response size for context budget"),
    ], size=12)

    _add_text(s, Inches(0.55), Inches(6.9), Inches(12.2), Inches(0.35),
              "Keyed by ANTHROPIC_API_KEY in backend .env; OpenAI provider stub exists for fallback.",
              size=11, color=MUTED)
    _footer(s, "AI / MCP")


def slide_integrations(prs):
    s = _slide(prs)
    _header(s, "Key External Integrations", "", "11")

    cards = [
        ("Cisco Cyber Vision", TEAL,
         "Primary downstream consumer — validation target.",
         ["API client queries CV device inventory",
          "MAC 100% / IP 95% confidence matching",
          "Enriches scenario devices with CV metadata",
          "Config: Settings → Cyber Vision"]),
        ("LDAP / Active Directory", NAVY,
         "Enterprise auth with local fallback.",
         ["ldap3 library; added Apr 2026",
          "Password validation via LDAP bind",
          "Group membership mapping",
          "Falls back to local users when LDAP unreachable"]),
        ("Anthropic Claude", ACCENT,
         "Scenario generation + in-product assistant.",
         ["Anthropic SDK 0.79 with MCP tools",
          "HTTP + Server-Sent Events transport",
          "Prompt caching supported",
          "OpenAI SDK present as fallback stub"]),
    ]

    x = Inches(0.55)
    for title, color, tag, items in cards:
        _add_rect(s, x, Inches(1.35), Inches(4.1), Inches(5.4), LIGHT)
        _add_rect(s, x, Inches(1.35), Inches(4.1), Inches(0.5), color)
        _add_text(s, x + Inches(0.2), Inches(1.42), Inches(3.8), Inches(0.4),
                  title, size=14, bold=True, color=WHITE)
        _add_text(s, x + Inches(0.2), Inches(1.95), Inches(3.8), Inches(0.5),
                  tag, size=11, bold=True, color=NAVY)
        _bullets(s, x + Inches(0.2), Inches(2.45), Inches(3.8), Inches(4.2),
                 items, size=11, color=SLATE, bold_lead=False)
        x += Inches(4.25)

    _footer(s, "Integrations")


def slide_dev(prs):
    s = _slide(prs)
    _header(s, "Development Environment", "Local loop — fast, reproducible, all-ports-on-host", "12")

    _add_text(s, Inches(0.55), Inches(1.25), Inches(6), Inches(0.4),
              "What Runs Where", size=14, bold=True, color=NAVY)
    _bullets(s, Inches(0.55), Inches(1.7), Inches(6), Inches(4), [
        ("Docker services:", "postgres:5432 · redis:6379 (via docker-compose.dev.yml)"),
        ("Backend on host:", "poetry run uvicorn --reload --port 8001"),
        ("Frontend on host:", "pnpm dev (Vite, :3001 or :5173)"),
        ("Services bind 0.0.0.0:", "off-box testing from other machines on the LAN"),
        ("Hot reload:", "uvicorn --reload + Vite HMR"),
    ], size=12)

    _add_text(s, Inches(6.9), Inches(1.25), Inches(6), Inches(0.4),
              "First-time Bootstrap", size=14, bold=True, color=NAVY)
    _bullets(s, Inches(6.9), Inches(1.7), Inches(6), Inches(4), [
        "cd backend && poetry lock && poetry install",
        "cd frontend && pnpm install",
        "cd docker && docker compose -f docker-compose.dev.yml up -d",
        "Start backend + frontend as above",
        "CORS allowed: localhost:3001, :5173",
    ], size=12)

    # port table
    _add_rect(s, Inches(0.55), Inches(5.5), Inches(12.2), Inches(1.5), LIGHT)
    _add_text(s, Inches(0.75), Inches(5.55), Inches(11.8), Inches(0.4),
              "Port Map", size=13, bold=True, color=NAVY)
    ports = [
        ("Backend", "8001"), ("Frontend dev", "3001 / 5173"),
        ("Postgres", "5432"), ("Redis", "6379"), ("pgAdmin (opt)", "5050"),
    ]
    x = Inches(0.75)
    for name, port in ports:
        _add_text(s, x, Inches(5.95), Inches(2.3), Inches(0.35),
                  name, size=12, bold=True, color=SLATE)
        _add_text(s, x, Inches(6.3), Inches(2.3), Inches(0.35),
                  port, size=14, bold=True, color=TEAL)
        x += Inches(2.4)

    _footer(s, "Dev Environment")


def slide_prod(prs):
    s = _slide(prs)
    _header(s, "Production Environment", "Same box, same compose file, different .env", "13")

    # row 1: nginx
    _add_rect(s, Inches(0.55), Inches(1.25), Inches(12.2), Inches(0.75), SLATE)
    _add_text(s, Inches(0.75), Inches(1.4), Inches(11.8), Inches(0.5),
              "Nginx (alpine)  ·  :443 HTTPS with self-signed TLS  ·  HTTP→HTTPS redirect  ·  WebSocket proxy  ·  SPA fallback",
              size=12, bold=True, color=WHITE)

    # row 2: backend / celery
    _add_rect(s, Inches(0.55), Inches(2.15), Inches(6.0), Inches(1.6), NAVY)
    _add_text(s, Inches(0.75), Inches(2.25), Inches(5.6), Inches(0.4),
              "backend (internal only)", size=13, bold=True, color=TEAL)
    _bullets(s, Inches(0.75), Inches(2.65), Inches(5.6), Inches(1.1), [
        "Uvicorn :8001, not exposed to host",
        "Mounts /var/run/docker.sock for agent image builds",
        "Healthcheck: GET /health",
    ], size=11, color=WHITE, bold_lead=False)

    _add_rect(s, Inches(6.75), Inches(2.15), Inches(6.0), Inches(1.6), NAVY)
    _add_text(s, Inches(6.95), Inches(2.25), Inches(5.6), Inches(0.4),
              "celery_worker", size=13, bold=True, color=TEAL)
    _bullets(s, Inches(6.95), Inches(2.65), Inches(5.6), Inches(1.1), [
        "Concurrency 2 — PCAP + deployment jobs",
        "Redis broker",
        "Healthcheck: celery inspect ping",
    ], size=11, color=WHITE, bold_lead=False)

    # row 3: data
    _add_rect(s, Inches(0.55), Inches(3.9), Inches(6.0), Inches(1.6), TEAL)
    _add_text(s, Inches(0.75), Inches(4.0), Inches(5.6), Inches(0.4),
              "postgres — TimescaleDB 15-alpine", size=13, bold=True, color=WHITE)
    _bullets(s, Inches(0.75), Inches(4.4), Inches(5.6), Inches(1.1), [
        "init-db.sql — uuid-ossp, pgcrypto, TimescaleDB",
        "Volume: postgres_data",
        "Healthcheck: pg_isready",
    ], size=11, color=WHITE, bold_lead=False)

    _add_rect(s, Inches(6.75), Inches(3.9), Inches(6.0), Inches(1.6), TEAL)
    _add_text(s, Inches(6.95), Inches(4.0), Inches(5.6), Inches(0.4),
              "redis — 7-alpine", size=13, bold=True, color=WHITE)
    _bullets(s, Inches(6.95), Inches(4.4), Inches(5.6), Inches(1.1), [
        "appendonly yes (AOF persistence)",
        "Volume: redis_data",
        "Celery broker + cache",
    ], size=11, color=WHITE, bold_lead=False)

    # row 4: secrets / deploy
    _add_rect(s, Inches(0.55), Inches(5.65), Inches(12.2), Inches(1.4), LIGHT)
    _add_text(s, Inches(0.75), Inches(5.72), Inches(11.8), Inches(0.4),
              "Secrets & Deploy", size=13, bold=True, color=NAVY)
    _bullets(s, Inches(0.75), Inches(6.1), Inches(11.8), Inches(1.0), [
        ("server-init.sh:", "generates POSTGRES_PASSWORD, SECRET_KEY, ADMIN_PASSWORD; chmod 600 .env"),
        ("SSL:", "auto-generated on first start; regen via volume reset"),
        ("Deploy cmd:", "docker compose up -d --build [service]  — same host, no orchestrator"),
    ], size=11, color=SLATE)

    _footer(s, "Prod Environment")


def slide_cicd(prs):
    s = _slide(prs)
    _header(s, "CI/CD", "GitHub Actions: lint → test → build → deploy", "14")

    # steps row
    steps = [
        ("1. Lint", "ruff (backend)\neslint + tsc (frontend)", NAVY),
        ("2. Test", "pytest (SQLite in-mem)\nvitest + jsdom + MSW", TEAL),
        ("3. Build", "docker compose build --no-cache\nup -d + /health check", ACCENT),
        ("4. Deploy", "SSH rsync → target host\ndocker compose up -d --build\nhealth loop 30×10s", SLATE),
    ]
    x = Inches(0.55)
    for title, body, color in steps:
        _add_rect(s, x, Inches(1.35), Inches(3.0), Inches(3.2), LIGHT)
        _add_rect(s, x, Inches(1.35), Inches(3.0), Inches(0.55), color)
        _add_text(s, x + Inches(0.15), Inches(1.43), Inches(2.8), Inches(0.4),
                  title, size=16, bold=True, color=WHITE)
        _add_text(s, x + Inches(0.15), Inches(2.1), Inches(2.8), Inches(2.3),
                  body, size=12, color=SLATE)
        x += Inches(3.1)

    _add_text(s, Inches(0.55), Inches(4.8), Inches(12.2), Inches(0.4),
              "Gates & Guardrails", size=14, bold=True, color=NAVY)
    _bullets(s, Inches(0.55), Inches(5.2), Inches(12.2), Inches(1.9), [
        "Backend coverage threshold: 50% (pytest-cov)",
        "CI uses SQLite in-memory DB — no external Postgres needed for unit tests",
        "Deploy via rsync over SSH; secrets from GitHub Actions (SSH_HOST, SSH_PRIVATE_KEY, POSTGRES_PASSWORD, SECRET_KEY, ADMIN_PASSWORD)",
        "Pre-commit hooks locally (ruff, black, mypy, type-check)",
        "Agent version bump is a code convention enforced by review, not CI — candidate for automation",
    ], size=12)

    _footer(s, "CI/CD")


def slide_risks(prs):
    s = _slide(prs)
    _header(s, "Strengths & Considerations", "Balanced view for planning", "15")

    _add_rect(s, Inches(0.55), Inches(1.3), Inches(6.0), Inches(5.6), LIGHT)
    _add_rect(s, Inches(0.55), Inches(1.3), Inches(6.0), Inches(0.55), TEAL)
    _add_text(s, Inches(0.75), Inches(1.38), Inches(5.6), Inches(0.4),
              "Strengths", size=16, bold=True, color=WHITE)
    _bullets(s, Inches(0.75), Inches(2.0), Inches(5.6), Inches(4.8), [
        "Modern, single-version stack across tiers",
        "Strong separation: canvas → orchestrator → engines → identity",
        "Shared protocol engines across backend and agent (one source of truth)",
        "Realism enforced by design: naming, protocols, MAC OUIs, IEC 62443 conduits",
        "AI is additive — product works fully without Claude",
        "Async end-to-end; TimescaleDB ready for time-series scale",
        "Agent phone-home model trivial to deploy behind customer firewalls",
    ], size=12)

    _add_rect(s, Inches(6.95), Inches(1.3), Inches(6.0), Inches(5.6), LIGHT)
    _add_rect(s, Inches(6.95), Inches(1.3), Inches(6.0), Inches(0.55), ACCENT)
    _add_text(s, Inches(7.15), Inches(1.38), Inches(5.6), Inches(0.4),
              "Considerations", size=16, bold=True, color=WHITE)
    _bullets(s, Inches(7.15), Inches(2.0), Inches(5.6), Inches(4.8), [
        "Single compose file for dev & prod — context drift risk; consider split compose + overrides",
        "Self-signed TLS in prod — fine for lab, blocker for enterprise; Let's Encrypt / cert mount path",
        "Monolith backend (25+ routers) — fine today, watch for slow test/startup over time",
        "Deploy = rsync + docker compose — no blue/green, no rollback target; consider image registry",
        "Agent version bump is a manual convention — add CI check",
        "Docker socket mounted into backend — capable but high-blast-radius; document threat model",
        "Secrets in .env on host — rotate story and backup story worth a slide of their own",
    ], size=12)

    _footer(s, "Strengths & Considerations")


def slide_closing(prs):
    s = _slide(prs)
    _add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    _add_rect(s, 0, Inches(3.2), SLIDE_W, Inches(0.05), TEAL)
    _add_text(s, Inches(0.75), Inches(2.4), Inches(12), Inches(0.8),
              "Summary", size=44, bold=True, color=WHITE)
    _add_text(s, Inches(0.75), Inches(3.4), Inches(12), Inches(0.5),
              "PacketArch is a coherent, modern, single-box platform.", size=20, color=TEAL)
    _bullets(s, Inches(0.75), Inches(4.1), Inches(12), Inches(2.8), [
        "Canvas-first UX · Async Python backend · 22+ OT protocols · 295 device templates",
        "Same engine code drives PCAP generation and live agent injection",
        "Realism is enforced, not optional — readiness gates every scenario",
        "Claude + MCP add natural-language scenario building without becoming load-bearing",
        "Prod = docker compose on a single host; next step is registry-based deploy + real TLS",
    ], size=14, color=WHITE, bold_lead=False)
    _add_text(s, Inches(0.75), Inches(6.9), Inches(12), Inches(0.4),
              "github.com/ip-aegis/PacketArch  ·  Architecture Review  ·  2026-04-24",
              size=11, color=MUTED)


# ---------------------------------------------------------------------------
def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs)
    slide_overview(prs)
    slide_system_arch(prs)
    slide_tech_stack(prs)
    slide_backend(prs)
    slide_frontend(prs)
    slide_agent(prs)
    slide_traffic_engine(prs)
    slide_protocols(prs)
    slide_fingerprints(prs)
    slide_ai(prs)
    slide_integrations(prs)
    slide_dev(prs)
    slide_prod(prs)
    slide_cicd(prs)
    slide_risks(prs)
    slide_closing(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    main()
