# PacketArch — OT Traffic Simulation Platform
# Copyright (c) 2026 Rocky Smith <rocky.d.smith@proton.me>
# Licensed under GPL-3.0. See LICENSE at the repo root.
"""Generate the Cisco Briefing PowerPoint for PacketArch.

Mirrors the 5-slide HTML deck at
backend/app/static/downloads/PacketArch-Cisco-Briefing.html
as a fully editable .pptx (every text box and rectangle is a native
PowerPoint shape, not a flattened image).

Run:
    python3 scripts/build_cisco_briefing_deck.py
Writes to:
    backend/app/static/downloads/PacketArch-Cisco-Briefing.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUT = (
    Path(__file__).resolve().parent.parent
    / "backend"
    / "app"
    / "static"
    / "downloads"
    / "PacketArch-Cisco-Briefing.pptx"
)

# --- Brand palette (matches the HTML deck) ---------------------------------
NAVY_DEEP = RGBColor(0x00, 0x1F, 0x2D)  # cover background, banners
NAVY_MID = RGBColor(0x00, 0x3A, 0x52)
NAVY = RGBColor(0x00, 0x50, 0x73)        # h2 strong accent
CYAN = RGBColor(0x00, 0xBC, 0xEB)        # primary accent
DARK_TEXT = RGBColor(0x00, 0x1F, 0x2D)
BODY = RGBColor(0x3A, 0x4A, 0x52)
MUTED = RGBColor(0x5A, 0x6A, 0x73)
CARD_BG = RGBColor(0xF5, 0xF8, 0xFA)
DIVIDER = RGBColor(0xE0, 0xE6, 0xEB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COVER_TAG = RGBColor(0xCF, 0xEA, 0xF3)
COVER_META = RGBColor(0x7F, 0xB8, 0xCC)
FOOTER_GREY = RGBColor(0x88, 0x88, 0x88)
BANNER_BODY = RGBColor(0xD8, 0xE6, 0xEC)

# --- Slide geometry --------------------------------------------------------
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
BRAND_BAR_H = Inches(0.08)
M_LEFT = Inches(0.7)
M_RIGHT = Inches(0.7)
TOP_EYEBROW = Inches(0.55)
CONTENT_W = SLIDE_W - M_LEFT - M_RIGHT
FOOTER_TOP = Inches(7.05)


# ---------- shape helpers --------------------------------------------------

def _rect(slide, left, top, width, height, fill, *, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def _text(
    slide,
    left,
    top,
    width,
    height,
    text,
    *,
    size=14,
    bold=False,
    color=BODY,
    align=PP_ALIGN.LEFT,
    font="Calibri",
    italic=False,
    anchor=MSO_ANCHOR.TOP,
    space_before=0,
):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    run.font.color.rgb = color
    return tb


def _rich(slide, left, top, width, height, segments, *, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """segments = [(text, size, bold, color, font), ...]"""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    for chunk, size, bold, color, font in segments:
        if not chunk:
            continue
        run = p.add_run()
        run.text = chunk
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = font
        run.font.color.rgb = color
    return tb


def _brand_bar(slide):
    _rect(slide, 0, 0, SLIDE_W, BRAND_BAR_H, CYAN)


def _footer(slide, page_num):
    _text(
        slide,
        M_LEFT,
        FOOTER_TOP,
        Inches(7),
        Inches(0.3),
        "PACKETARCH  ·  CISCO BRIEFING",
        size=10,
        color=FOOTER_GREY,
    )
    _text(
        slide,
        SLIDE_W - Inches(2.2) - M_RIGHT,
        FOOTER_TOP,
        Inches(2.2),
        Inches(0.3),
        page_num,
        size=10,
        bold=True,
        color=NAVY,
        align=PP_ALIGN.RIGHT,
    )


def _eyebrow(slide, text, top=TOP_EYEBROW):
    _text(
        slide,
        M_LEFT,
        top,
        CONTENT_W,
        Inches(0.3),
        text.upper(),
        size=11,
        bold=True,
        color=CYAN,
    )


def _h2(slide, top, segments):
    """segments = [(text, bold), ...] — light Calibri Light by default,
    bolded chunks use Calibri + NAVY for emphasis."""
    full_segments = []
    for txt, bold in segments:
        font = "Calibri" if bold else "Calibri Light"
        color = NAVY if bold else DARK_TEXT
        full_segments.append((txt, 36, bold, color, font))
    _rich(slide, M_LEFT, top, CONTENT_W, Inches(0.9), full_segments)


def _subhead(slide, top, text, width=None):
    _text(
        slide,
        M_LEFT,
        top,
        width or CONTENT_W,
        Inches(0.7),
        text,
        size=16,
        color=MUTED,
        font="Calibri Light",
    )


# ---------- slide builders -------------------------------------------------

def slide_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY_DEEP)
    _rect(s, 0, 0, SLIDE_W, BRAND_BAR_H, CYAN)

    _text(
        s,
        Inches(1.0),
        Inches(1.7),
        Inches(11),
        Inches(0.4),
        "OT  TRAFFIC  SIMULATION  PLATFORM",
        size=14,
        bold=True,
        color=CYAN,
    )
    # Title — "PacketArch" with bold "Arch"
    _rich(
        s,
        Inches(1.0),
        Inches(2.2),
        Inches(11),
        Inches(1.6),
        [
            ("Packet", 88, False, WHITE, "Calibri Light"),
            ("Arch", 88, True, WHITE, "Calibri"),
        ],
    )
    _text(
        s,
        Inches(1.0),
        Inches(3.9),
        Inches(11),
        Inches(1.2),
        "Fingerprint-grade industrial traffic, generated on demand — "
        "built with Cisco Cyber Vision in mind.",
        size=24,
        color=COVER_TAG,
        font="Calibri Light",
    )
    _text(
        s,
        Inches(1.0),
        Inches(5.5),
        Inches(11),
        Inches(0.4),
        "CISCO BRIEFING   ·   MAY 2026",
        size=12,
        color=COVER_META,
    )


def slide_gap(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _brand_bar(s)
    _eyebrow(s, "The Gap")
    _h2(
        s,
        Inches(0.85),
        [
            ("You can't ", False),
            ("demo, train, or validate", True),
            (" on a live plant.", False),
        ],
    )
    _subhead(
        s,
        Inches(1.55),
        "OT networks are quiet by design. The traffic Cyber Vision needs to "
        "classify a plant only exists inside the plant — until now.",
    )

    # Three cards
    cards = [
        (
            "01",
            "Production is off-limits",
            "Safety, change control, and uptime mean you can't introduce test "
            "traffic into a running OT environment.",
        ),
        (
            "02",
            "Lab PCAPs don't scale",
            "Captures are frozen, narrow, and don't adapt. A handful of PCAPs "
            "can't represent the diversity of a real plant.",
        ),
        (
            "03",
            "Generators look fake",
            "Generic packet tools don't carry the vendor OUI, sysObjectID, or "
            "protocol identity that Cyber Vision keys on.",
        ),
    ]
    card_top = Inches(2.7)
    card_h = Inches(2.05)
    gap = Inches(0.25)
    card_w = (CONTENT_W - gap * 2) / 3
    for i, (num, title, body) in enumerate(cards):
        left = M_LEFT + (card_w + gap) * i
        _rect(s, left, card_top, card_w, card_h, CARD_BG)
        # Cyan left border
        _rect(s, left, card_top, Inches(0.05), card_h, CYAN)
        _text(
            s,
            left + Inches(0.25),
            card_top + Inches(0.2),
            card_w - Inches(0.5),
            Inches(0.3),
            num,
            size=11,
            bold=True,
            color=CYAN,
        )
        _text(
            s,
            left + Inches(0.25),
            card_top + Inches(0.55),
            card_w - Inches(0.5),
            Inches(0.5),
            title,
            size=18,
            bold=True,
            color=DARK_TEXT,
        )
        _text(
            s,
            left + Inches(0.25),
            card_top + Inches(1.05),
            card_w - Inches(0.5),
            Inches(0.95),
            body,
            size=13,
            color=BODY,
        )

    # Banner
    banner_top = Inches(5.05)
    banner_h = Inches(1.0)
    _rect(s, M_LEFT, banner_top, CONTENT_W, banner_h, NAVY_DEEP)
    _rich(
        s,
        M_LEFT + Inches(0.4),
        banner_top + Inches(0.18),
        CONTENT_W - Inches(0.8),
        banner_h - Inches(0.3),
        [
            ("The unmet need:  ", 18, False, WHITE, "Calibri Light"),
            (
                "on-demand, fingerprint-grade OT traffic",
                18,
                True,
                CYAN,
                "Calibri",
            ),
            (
                "  that Cyber Vision treats as the real thing.",
                18,
                False,
                WHITE,
                "Calibri Light",
            ),
        ],
        anchor=MSO_ANCHOR.MIDDLE,
    )

    _footer(s, "02 / 05")


def slide_platform(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _brand_bar(s)
    _eyebrow(s, "The Platform")
    _h2(
        s,
        Inches(0.85),
        [("A ", False), ("scenario studio", True), (" for industrial networks.", False)],
    )
    _subhead(
        s,
        Inches(1.55),
        "Design a plant visually. Generate traffic that looks like that plant. "
        "Run it offline as PCAP or live through remote agents.",
    )

    caps = [
        ("D", "Visual Scenario Design",
         "Purdue-aware canvas, IEC 62443 conduits, drag-and-drop devices and zones."),
        ("F", "Vendor Fingerprint Library",
         "295 device templates across 18 OT vendors — Siemens, Rockwell, Schneider, Honeywell, ABB, Yokogawa, Cisco, and more."),
        ("P", "Six Production Protocols",
         "Modbus TCP, EtherNet/IP, PROFINET, S7comm, BACnet/IP, SNMP/NTCIP — with full identity payloads."),
        ("L", "Live Agents or PCAP",
         "Remote agents inject real packets on a SPAN port, or generate offline PCAPs for analysis and replay."),
        ("A", "AI Scenario Generation",
         "Natural-language to full plant in seconds. Backed by Claude, OpenAI, or Cisco's CIRCUIT gateway."),
        ("X", "Attack & Adaptive Engines",
         "TRITON / PIPEDREAM / INDUSTROYER playbooks, MITRE ATT&CK for ICS mapping, time-of-day traffic drift."),
    ]
    grid_top = Inches(2.55)
    row_h = Inches(1.15)
    col_gap = Inches(0.5)
    cell_w = (CONTENT_W - col_gap) / 2
    icon_size = Inches(0.5)
    for i, (icon, title, body) in enumerate(caps):
        col = i % 2
        row = i // 2
        left = M_LEFT + (cell_w + col_gap) * col
        top = grid_top + row_h * row
        _rect(s, left, top, icon_size, icon_size, CYAN)
        _text(
            s,
            left,
            top,
            icon_size,
            icon_size,
            icon,
            size=20,
            bold=True,
            color=WHITE,
            align=PP_ALIGN.CENTER,
            font="Cambria",
            anchor=MSO_ANCHOR.MIDDLE,
        )
        body_left = left + icon_size + Inches(0.18)
        body_w = cell_w - icon_size - Inches(0.18)
        _text(
            s,
            body_left,
            top - Inches(0.02),
            body_w,
            Inches(0.4),
            title,
            size=16,
            bold=True,
            color=DARK_TEXT,
        )
        _text(
            s,
            body_left,
            top + Inches(0.35),
            body_w,
            Inches(0.8),
            body,
            size=12,
            color=BODY,
        )

    # Stat strip
    strip_top = Inches(6.05)
    _rect(s, M_LEFT, strip_top, CONTENT_W, Inches(0.02), DIVIDER)
    stats = [
        ("295", "DEVICE FINGERPRINTS"),
        ("18", "OT VENDORS"),
        ("6", "PROTOCOLS"),
        ("6", "INDUSTRY VERTICALS"),
        ("9", "ATTACK PLAYBOOKS"),
    ]
    stat_w = CONTENT_W / 5
    for i, (num, label) in enumerate(stats):
        left = M_LEFT + stat_w * i
        _text(
            s,
            left,
            strip_top + Inches(0.15),
            stat_w,
            Inches(0.55),
            num,
            size=32,
            color=NAVY,
            align=PP_ALIGN.CENTER,
            font="Calibri Light",
        )
        _text(
            s,
            left,
            strip_top + Inches(0.7),
            stat_w,
            Inches(0.25),
            label,
            size=9,
            bold=True,
            color=MUTED,
            align=PP_ALIGN.CENTER,
        )

    _footer(s, "03 / 05")


def slide_workflow(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _brand_bar(s)
    _eyebrow(s, "The Workflow")
    _h2(
        s,
        Inches(0.85),
        [("From ", False), ("idea to live demo", True), (" in minutes.", False)],
    )
    _subhead(
        s,
        Inches(1.55),
        "One sentence into the AI Wizard becomes a complete plant, a "
        "fingerprint-grade PCAP, and a Cyber Vision demo ready to run.",
    )

    # 4 step cards with 3 chevron arrows between
    steps = [
        ("STEP 01", "Describe",
         "One sentence into the AI Wizard — the plant, the vendors, the rough size."),
        ("STEP 02", "Generate",
         "Full scenario auto-built — devices, protocols, conduits, fingerprints, flows."),
        ("STEP 03", "Export PCAP",
         "One click. Realistic, time-stamped industrial traffic written to a single file."),
        ("STEP 04", "Import to CV",
         "Drop the PCAP into Cyber Vision. Every device classifies on first look. Demo runs."),
    ]
    pipe_top = Inches(2.55)
    pipe_h = Inches(1.85)
    arrow_w = Inches(0.35)
    total_arrows_w = arrow_w * 3
    card_w = (CONTENT_W - total_arrows_w) / 4
    x = M_LEFT
    for i, (label, title, body) in enumerate(steps):
        # Card body
        _rect(s, x, pipe_top, card_w, pipe_h, CARD_BG)
        # Cyan top border
        _rect(s, x, pipe_top, card_w, Inches(0.06), CYAN)
        _text(
            s,
            x + Inches(0.22),
            pipe_top + Inches(0.18),
            card_w - Inches(0.44),
            Inches(0.3),
            label,
            size=10,
            bold=True,
            color=CYAN,
        )
        _text(
            s,
            x + Inches(0.22),
            pipe_top + Inches(0.5),
            card_w - Inches(0.44),
            Inches(0.5),
            title,
            size=20,
            bold=True,
            color=DARK_TEXT,
        )
        _text(
            s,
            x + Inches(0.22),
            pipe_top + Inches(1.0),
            card_w - Inches(0.44),
            Inches(0.85),
            body,
            size=12,
            color=BODY,
        )
        x += card_w
        if i < len(steps) - 1:
            # Chevron between steps
            _text(
                s,
                x,
                pipe_top,
                arrow_w,
                pipe_h,
                "›",
                size=40,
                bold=True,
                color=CYAN,
                align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE,
            )
            x += arrow_w

    # Example callout
    ex_top = Inches(4.7)
    ex_h = Inches(2.2)
    _rect(s, M_LEFT, ex_top, CONTENT_W, ex_h, NAVY_DEEP)
    _rect(s, M_LEFT, ex_top, Inches(0.06), ex_h, CYAN)
    _text(
        s,
        M_LEFT + Inches(0.35),
        ex_top + Inches(0.2),
        CONTENT_W - Inches(0.7),
        Inches(0.3),
        "EXAMPLE  ·  THE MONDAY-MORNING BAKERY DEMO",
        size=11,
        bold=True,
        color=CYAN,
    )
    # Rich text story
    tb = s.shapes.add_textbox(
        M_LEFT + Inches(0.35),
        ex_top + Inches(0.55),
        CONTENT_W - Inches(0.7),
        ex_h - Inches(0.7),
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    story = [
        ("Customer asks for a Cyber Vision walkthrough on Monday. You type one sentence into the AI Wizard: ", False),
        ("“commercial bakery — 6 ovens, 2 industrial mixers, a packaging line, and building automation.”", True),
        (" Minutes later you have a complete plant — Siemens S7-1500 oven controllers on PROFINET, Rockwell ControlLogix on the mixers, BACnet thermostats and AHUs for the BMS, all wired with realistic flow patterns. Click ", False),
        ("Export PCAP", True),
        (". Drop the file into Cyber Vision. Every device classifies on first look. Demo done over coffee.", False),
    ]
    for chunk, bold in story:
        run = p.add_run()
        run.text = chunk
        run.font.size = Pt(14)
        run.font.bold = bold
        run.font.name = "Calibri"
        run.font.color.rgb = WHITE if bold else BANNER_BODY

    _footer(s, "04 / 05")


def slide_cisco_value(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _brand_bar(s)
    _eyebrow(s, "Why Cisco")
    _h2(
        s,
        Inches(0.85),
        [("Built ", False), ("with Cyber Vision", True), (" in mind.", False)],
    )
    _subhead(
        s,
        Inches(1.55),
        "Every packet is engineered to be classified correctly — the same "
        "fields, OIDs, and identifiers CV uses in production.",
    )

    # Two columns
    col_top = Inches(2.65)
    col_gap = Inches(0.5)
    col_w = (CONTENT_W - col_gap) / 2
    col1_left = M_LEFT
    col2_left = M_LEFT + col_w + col_gap

    _text(s, col1_left, col_top, col_w, Inches(0.35),
          "CYBER VISION NATIVE", size=14, bold=True, color=NAVY)
    _rect(s, col1_left, col_top + Inches(0.42), col_w, Inches(0.025), CYAN)

    _text(s, col2_left, col_top, col_w, Inches(0.35),
          "WHERE IT PAYS OFF", size=14, bold=True, color=NAVY)
    _rect(s, col2_left, col_top + Inches(0.42), col_w, Inches(0.025), CYAN)

    cv_items = [
        ("Direct CV integration",
         " — query the center, match devices, see fingerprint deltas inline."),
        ("Fingerprint-grade payloads",
         " — correct vendor OUI, sysObjectID, Modbus MEI, S7 SZL, BACnet Object_Identifier, CIP Identity Object."),
        ("CIRCUIT AI provider",
         " — scenario generation runs on Cisco's internal LLM gateway, no external API exposure."),
        ("Air-gapped lab kits",
         " — offline tarball bundles ship with everything for isolated or secure labs."),
    ]
    use_items = [
        ("Pre-sales demos",
         " — stand up a 200-device water plant in minutes and let CV classify it live in front of the customer."),
        ("Internal & partner training",
         " — hands-on labs for Cisco SEs, TAC, and partners on realistic OT environments — no production access required."),
        ("Field enablement",
         " — every SE carries a complete OT plant in a laptop, ready to demo any vertical on demand."),
    ]

    def _bullet_block(left, top, width, items, dot_color):
        item_top = top
        for lead, rest in items:
            # Dot
            _rect(
                s,
                left + Inches(0.05),
                item_top + Inches(0.13),
                Inches(0.1),
                Inches(0.1),
                dot_color,
            )
            tb = s.shapes.add_textbox(left + Inches(0.3), item_top, width - Inches(0.3), Inches(0.9))
            tf = tb.text_frame
            tf.word_wrap = True
            tf.margin_left = Inches(0.02)
            tf.margin_right = Inches(0.02)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            r1 = p.add_run()
            r1.text = lead
            r1.font.size = Pt(13)
            r1.font.bold = True
            r1.font.name = "Calibri"
            r1.font.color.rgb = DARK_TEXT
            r2 = p.add_run()
            r2.text = rest
            r2.font.size = Pt(13)
            r2.font.name = "Calibri"
            r2.font.color.rgb = BODY
            item_top += Inches(0.9)

    _bullet_block(col1_left, col_top + Inches(0.65), col_w, cv_items, CYAN)

    # Right column has a subtle background panel
    panel_top = col_top + Inches(0.65)
    panel_h = Inches(0.9) * len(use_items) + Inches(0.2)
    _rect(s, col2_left, panel_top, col_w, panel_h, CARD_BG)
    _bullet_block(col2_left + Inches(0.15), panel_top + Inches(0.1), col_w - Inches(0.3), use_items, NAVY)

    _footer(s, "05 / 05")


# ---------- main -----------------------------------------------------------

def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_cover(prs)
    slide_gap(prs)
    slide_platform(prs)
    slide_workflow(prs)
    slide_cisco_value(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"Wrote {OUT}  ({OUT.stat().st_size} bytes)")


if __name__ == "__main__":
    build()
